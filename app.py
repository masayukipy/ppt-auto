import collections  # noqa
import collections.abc  # noqa
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from utils.scrape import Codephil

font_size = 42

Path("ppt").mkdir(parents=True, exist_ok=True)


def create_slide(prs, chapter):
    text = chapter["name"]
    mul_text = ""
    while True:
        if len(text) > 15:
            mul_text = text[:15] + "\n"
            text = text[15:]
        else:
            mul_text += text
            break

    mul_len = len(mul_text.split("\n"))

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # text box
    tx_box = slide.shapes.add_textbox(0, 0, prs.slide_width, Pt(font_size) * mul_len)
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = mul_text
    p.alignment = PP_ALIGN.CENTER

    # font
    font = run.font
    font.name = "MS Mincho"
    font.size = Pt(font_size)
    tx_box.left = int((prs.slide_width - tx_box.width) / 2)
    tx_box.top = int(
        (prs.slide_height - Pt(font_size + 20 + (mul_len - 1) * font_size)) / 2
    )

    # image
    if not chapter["img_path"] is None:
        slide.shapes.add_picture(
            chapter["img_path"], 0, 0, prs.slide_width, prs.slide_height
        )


def all():
    article_list = Codephil().run()
    for article in article_list:
        prs = Presentation()
        for chapter in article:
            create_slide(prs, chapter)
        prs.save(f'ppt/{article[0]["name"]}.pptx')


all()
