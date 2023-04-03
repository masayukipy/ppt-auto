from PIL import Image
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

font_size = 52
col_letters = 12


def create_slide(prs, chapter):
    prs.slide_height = int(prs.slide_width * 9 / 16)
    text = chapter["name"]
    mul_text = ""
    while True:
        if len(text) > col_letters:
            mul_text += text[:col_letters] + "\n"
            text = text[col_letters:]
        else:
            mul_text += text
            break

    mul_len = len(mul_text.split("\n"))

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # text box
    tx_box = slide.shapes.add_textbox(0, 0, prs.slide_width, Pt(font_size) * mul_len)
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = mul_text
    p.alignment = PP_ALIGN.CENTER

    # font
    font = run.font
    font.name = "Arial"
    font.size = Pt(font_size)
    tx_box.left = int((prs.slide_width - tx_box.width) / 2)
    tx_box.top = int(
        (prs.slide_height - Pt(font_size + 20 + (mul_len - 1) * font_size)) / 2
    )

    # image
    if not chapter["img_path"] is None:
        img = Image.open(chapter["img_path"])
        hg = int(img.height * prs.slide_width / img.width)
        slide.shapes.add_picture(
            chapter["img_path"],
            0,
            int(0 - (hg - prs.slide_height) / 2),
            prs.slide_width,
            hg,
        )
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = chapter["content"]
    return
